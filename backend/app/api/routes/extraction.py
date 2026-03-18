"""
Additional API routes for voice transcription, document extraction, task parsing, and LinkedIn extraction
"""
from fastapi import APIRouter, File, UploadFile, HTTPException
from pydantic import BaseModel
from typing import List, Dict, Optional
import os
from anthropic import Anthropic
import pypdf
import docx
import tempfile
import base64
import csv
import json
import io
import re
import urllib.request
import urllib.parse

router = APIRouter()


class ParseTasksRequest(BaseModel):
    text: str


class TaskExtract(BaseModel):
    name: str
    description: str
    frequency: str
    time_per_task: int
    category: str
    complexity: str


class ParsedTasksResponse(BaseModel):
    workflow_name: str
    workflow_description: str
    tasks: List[TaskExtract]


# All supported image media types for Claude Vision
IMAGE_MEDIA_TYPES = {
    'png': 'image/png',
    'jpg': 'image/jpeg',
    'jpeg': 'image/jpeg',
    'gif': 'image/gif',
    'webp': 'image/webp',
    'bmp': 'image/png',    # convert via re-encode
    'tiff': 'image/png',
    'tif': 'image/png',
    'heic': 'image/jpeg',
    'heif': 'image/jpeg',
    'avif': 'image/png',
    'svg': 'image/png',
    'ico': 'image/png',
}

# All supported text/document formats
TEXT_EXTENSIONS = {
    'txt', 'md', 'markdown', 'rst', 'rtf',
    'csv', 'tsv',
    'json', 'jsonl',
    'xml', 'html', 'htm',
    'yaml', 'yml',
    'log',
}


def extract_text_from_file(file_path: str, filename: str) -> str:
    """Extract text from 20+ document and image formats"""
    ext = filename.lower().rsplit('.', 1)[-1] if '.' in filename else ''

    try:
        # ── Plain text variants ────────────────────────────────────────────
        if ext in TEXT_EXTENSIONS:
            if ext in ('csv', 'tsv'):
                delimiter = '\t' if ext == 'tsv' else ','
                with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                    reader = csv.reader(f, delimiter=delimiter)
                    rows = list(reader)
                lines = []
                if rows:
                    lines.append(' | '.join(str(c) for c in rows[0]))  # header
                    for row in rows[1:]:
                        lines.append(' | '.join(str(c) for c in row))
                return '\n'.join(lines)

            if ext == 'json':
                with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                    data = json.load(f)
                return json.dumps(data, indent=2, ensure_ascii=False)

            with open(file_path, 'r', encoding='utf-8', errors='replace') as f:
                return f.read()

        # ── PDF ───────────────────────────────────────────────────────────
        elif ext == 'pdf':
            text = ""
            with open(file_path, 'rb') as f:
                pdf_reader = pypdf.PdfReader(f)
                for page in pdf_reader.pages:
                    text += (page.extract_text() or '') + "\n"
            return text.strip() or "[PDF had no extractable text — may be scanned]"

        # ── Word documents ────────────────────────────────────────────────
        elif ext in ('doc', 'docx'):
            doc = docx.Document(file_path)
            parts = []
            for para in doc.paragraphs:
                if para.text.strip():
                    parts.append(para.text)
            for table in doc.tables:
                for row in table.rows:
                    parts.append(' | '.join(c.text.strip() for c in row.cells if c.text.strip()))
            return '\n'.join(parts)

        # ── ODT (LibreOffice) ─────────────────────────────────────────────
        elif ext == 'odt':
            try:
                from odf.opendocument import load as odf_load
                from odf.text import P
                from odf.element import Element
                doc = odf_load(file_path)
                texts = []
                for elem in doc.text.getElementsByType(P):
                    t = elem.plaintext() if hasattr(elem, 'plaintext') else str(elem)
                    if t.strip():
                        texts.append(t.strip())
                return '\n'.join(texts) if texts else "[ODT: no text extracted]"
            except ImportError:
                # Fallback: read raw XML
                import zipfile, re as _re
                with zipfile.ZipFile(file_path) as z:
                    with z.open('content.xml') as f:
                        raw = f.read().decode('utf-8', errors='replace')
                return _re.sub(r'<[^>]+>', ' ', raw)

        # ── Excel ─────────────────────────────────────────────────────────
        elif ext in ('xls', 'xlsx', 'xlsm', 'xlsb', 'ods'):
            try:
                import openpyxl
                wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
                parts = []
                for sheet_name in wb.sheetnames:
                    ws = wb[sheet_name]
                    parts.append(f"=== Sheet: {sheet_name} ===")
                    for row in ws.iter_rows(values_only=True):
                        if any(c is not None for c in row):
                            parts.append(' | '.join(str(c) if c is not None else '' for c in row))
                return '\n'.join(parts)
            except Exception:
                # Try xlrd for older xls
                try:
                    import xlrd
                    wb = xlrd.open_workbook(file_path)
                    parts = []
                    for sheet in wb.sheets():
                        parts.append(f"=== Sheet: {sheet.name} ===")
                        for rx in range(sheet.nrows):
                            parts.append(' | '.join(str(sheet.cell_value(rx, cx)) for cx in range(sheet.ncols)))
                    return '\n'.join(parts)
                except Exception as e:
                    raise ValueError(f"Could not read spreadsheet: {e}")

        # ── PowerPoint ────────────────────────────────────────────────────
        elif ext in ('ppt', 'pptx'):
            try:
                from pptx import Presentation
                prs = Presentation(file_path)
                parts = []
                for i, slide in enumerate(prs.slides, 1):
                    parts.append(f"=== Slide {i} ===")
                    for shape in slide.shapes:
                        if hasattr(shape, 'text') and shape.text.strip():
                            parts.append(shape.text.strip())
                return '\n'.join(parts)
            except ImportError:
                raise ValueError("pptx support requires python-pptx — install it on the server")

        # ── Images — all formats via Claude Vision ────────────────────────
        elif ext in IMAGE_MEDIA_TYPES:
            api_key = os.getenv("ANTHROPIC_API_KEY")
            if not api_key:
                raise ValueError("ANTHROPIC_API_KEY not configured for image OCR")

            # For formats Claude Vision doesn't natively support, convert to PNG first
            native_formats = {'png', 'jpg', 'jpeg', 'gif', 'webp'}
            if ext not in native_formats:
                try:
                    from PIL import Image as PILImage
                    img = PILImage.open(file_path).convert('RGB')
                    converted_path = file_path + '_converted.png'
                    img.save(converted_path, 'PNG')
                    read_path = converted_path
                    media_type = 'image/png'
                except ImportError:
                    read_path = file_path
                    media_type = IMAGE_MEDIA_TYPES.get(ext, 'image/png')
            else:
                read_path = file_path
                media_type = IMAGE_MEDIA_TYPES[ext]

            with open(read_path, 'rb') as f:
                image_data = base64.b64encode(f.read()).decode('utf-8')

            client = Anthropic(api_key=api_key)
            message = client.messages.create(
                model="claude-sonnet-4-20250514",
                max_tokens=2000,
                messages=[{
                    "role": "user",
                    "content": [
                        {
                            "type": "image",
                            "source": {"type": "base64", "media_type": media_type, "data": image_data}
                        },
                        {
                            "type": "text",
                            "text": "Extract ALL text from this image. Return ONLY the extracted text, preserving structure. If this is a workflow, task list, or process diagram, preserve all labels and steps."
                        }
                    ]
                }]
            )
            return message.content[0].text

        else:
            raise ValueError(f"Unsupported file type: .{ext}")

    except HTTPException:
        raise
    except Exception as e:
        raise HTTPException(status_code=400, detail=f"Failed to extract text: {str(e)}")


class LinkedInRequest(BaseModel):
    url: str
    profile_type: Optional[str] = None  # 'personal' or 'company'


class LinkedInExtractResponse(BaseModel):
    text: str
    profile_type: str  # 'personal' or 'company'
    name: str
    title_or_tagline: str


def _detect_linkedin_type(url: str) -> str:
    """Detect whether URL is a personal profile or company page."""
    clean = url.lower().strip().rstrip('/')
    if '/company/' in clean or '/school/' in clean or '/showcase/' in clean:
        return 'company'
    return 'personal'


def _fetch_linkedin_html(url: str) -> str:
    """
    Fetch LinkedIn page HTML server-side.
    Uses a realistic browser User-Agent to get rendered meta/og tags and basic page structure.
    LinkedIn returns a lot of JS-gated content, but the og: meta tags + early HTML
    contain enough name/title/about/experience for our purposes.
    """
    headers = {
        'User-Agent': (
            'Mozilla/5.0 (Windows NT 10.0; Win64; x64) '
            'AppleWebKit/537.36 (KHTML, like Gecko) '
            'Chrome/124.0.0.0 Safari/537.36'
        ),
        'Accept': 'text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8',
        'Accept-Language': 'en-US,en;q=0.9',
        'Accept-Encoding': 'gzip, deflate',
        'Connection': 'keep-alive',
        'Upgrade-Insecure-Requests': '1',
    }
    req = urllib.request.Request(url, headers=headers)
    try:
        with urllib.request.urlopen(req, timeout=12) as response:
            raw = response.read()
            # Handle gzip
            encoding = response.headers.get('Content-Encoding', '')
            if encoding == 'gzip':
                import gzip
                raw = gzip.decompress(raw)
            return raw.decode('utf-8', errors='replace')
    except Exception as e:
        raise HTTPException(status_code=422, detail=f"Could not fetch LinkedIn page: {str(e)}")


def _extract_linkedin_signals(html: str, profile_type: str) -> str:
    """
    Pull the most signal-rich text from LinkedIn HTML without a full HTML parser.
    Focuses on: og: meta tags, title, JSON-LD structured data, visible text snippets.
    """
    signals = []

    # og: meta tags (most reliable — always present even with JS gating)
    og_tags = re.findall(r'<meta[^>]+property=["\']og:([^"\']+)["\'][^>]+content=["\']([^"\']{1,600})["\']', html, re.I)
    og_tags += re.findall(r'<meta[^>]+content=["\']([^"\']{1,600})["\'][^>]+property=["\']og:([^"\']+)["\']', html, re.I)
    for tag in og_tags:
        key = tag[0] if ':' not in tag[0] else tag[0]
        val = tag[1]
        signals.append(f"og:{key}: {val}")

    # Page title
    title_m = re.search(r'<title[^>]*>([^<]{1,300})</title>', html, re.I)
    if title_m:
        signals.append(f"page_title: {title_m.group(1).strip()}")

    # JSON-LD structured data
    ld_blocks = re.findall(r'<script[^>]+type=["\']application/ld\+json["\'][^>]*>(.*?)</script>', html, re.I | re.S)
    for block in ld_blocks[:3]:
        try:
            data = json.loads(block)
            signals.append(f"structured_data: {json.dumps(data, ensure_ascii=False)[:1500]}")
        except Exception:
            pass

    # Inline JSON data blobs (LinkedIn puts profile data in window.__initialData__ or similar)
    json_blobs = re.findall(r'(?:voyagerIdentity|miniProfile|firstName|lastName|headline|summary|description|companyPageUrl)["\']?\s*:\s*["\']([^"\'<]{1,500})', html)
    for blob in json_blobs[:20]:
        signals.append(blob)

    # Visible text from key semantic areas — strip tags and collapse whitespace
    body_text = re.sub(r'<script[^>]*>.*?</script>', ' ', html, flags=re.S | re.I)
    body_text = re.sub(r'<style[^>]*>.*?</style>', ' ', body_text, flags=re.S | re.I)
    body_text = re.sub(r'<[^>]+>', ' ', body_text)
    body_text = re.sub(r'\s+', ' ', body_text).strip()
    # Take first 3000 chars of visible text — usually contains the above-the-fold profile info
    signals.append(f"visible_text_excerpt: {body_text[:3000]}")

    return '\n'.join(signals)


@router.post("/extract-linkedin", response_model=LinkedInExtractResponse)
async def extract_linkedin_profile(request: LinkedInRequest):
    """
    Fetch a LinkedIn personal profile or company page and extract role/responsibilities context.
    Uses server-side fetch + Claude Haiku for text extraction and structuring.
    No external services — uses Anthropic API already integrated.
    """
    url = request.url.strip()

    # Basic URL validation
    if not url.startswith('http'):
        url = 'https://' + url
    if 'linkedin.com' not in url.lower():
        raise HTTPException(status_code=400, detail="Please provide a valid LinkedIn URL (linkedin.com/in/... or linkedin.com/company/...)")

    profile_type = request.profile_type or _detect_linkedin_type(url)

    api_key = os.getenv("ANTHROPIC_API_KEY")
    if not api_key:
        raise HTTPException(status_code=500, detail="ANTHROPIC_API_KEY not configured")

    # Fetch the page
    html = _fetch_linkedin_html(url)
    signals = _extract_linkedin_signals(html, profile_type)

    client = Anthropic(api_key=api_key)

    if profile_type == 'personal':
        system_prompt = """You are an expert at reading LinkedIn profile data from raw HTML signals and meta tags.
Extract the person's professional context to enable an AI workflow automation analysis of their role and daily responsibilities.
Be thorough and infer tasks even when not explicitly listed — use title, industry, company size as context."""
        user_prompt = f"""LinkedIn personal profile signals extracted from HTML:

{signals[:4000]}

From these signals, produce a detailed professional context in plain text (not JSON) describing:
1. Person's full name and current job title / role
2. Current employer and industry
3. Key responsibilities and daily tasks (infer from title/industry if not explicit)
4. Notable skills and tools used
5. Career level and specialisation

Be specific and detailed — this text will be used to generate an AI automation readiness analysis.
Write it as a professional profile summary (3-6 paragraphs)."""
    else:
        system_prompt = """You are an expert at reading LinkedIn company page data from raw HTML signals and meta tags.
Extract the company's context to enable an AI workflow automation analysis of their typical workflows, team structure and daily operations."""
        user_prompt = f"""LinkedIn company page signals extracted from HTML:

{signals[:4000]}

From these signals, produce a detailed company context in plain text (not JSON) describing:
1. Company name, industry, and size
2. Core business activities and what teams do day-to-day
3. Key departments likely present and their typical tasks
4. Products/services offered and operational processes involved
5. Technology stack or tools mentioned

Be specific and detailed — this text will be used to generate an AI automation readiness analysis.
Write it as a company profile summary (3-6 paragraphs)."""

    try:
        message = client.messages.create(
            model="claude-haiku-4-5-20251001",
            max_tokens=1200,
            system=system_prompt,
            messages=[{"role": "user", "content": user_prompt}]
        )
        extracted_text = message.content[0].text.strip()
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"AI extraction failed: {str(e)}")

    # Pull name and title/tagline from og tags for display
    name_match = re.search(r'og:title:\s*(.+?)(?:\s*[-|]\s*LinkedIn|\n|$)', signals, re.I)
    desc_match = re.search(r'og:description:\s*(.+?)(?:\n|$)', signals, re.I)
    name = name_match.group(1).strip() if name_match else ('Profile' if profile_type == 'personal' else 'Company')
    title_or_tagline = desc_match.group(1).strip()[:120] if desc_match else ''

    return LinkedInExtractResponse(
        text=extracted_text,
        profile_type=profile_type,
        name=name,
        title_or_tagline=title_or_tagline,
    )


@router.post("/transcribe")
async def transcribe_audio(audio: UploadFile = File(...)):
    return {"transcription": "Please describe your workflow tasks manually for now."}


@router.post("/extract-tasks")
async def extract_tasks_from_document(file: UploadFile = File(...)):
    """Extract text from uploaded document — supports 20+ formats"""
    ext = (file.filename or '').rsplit('.', 1)[-1].lower()
    suffix = f".{ext}" if ext else ""

    with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
        content = await file.read()
        tmp.write(content)
        tmp_path = tmp.name

    try:
        text = extract_text_from_file(tmp_path, file.filename or "upload")
        return {"text": text}
    finally:
        if os.path.exists(tmp_path):
            os.remove(tmp_path)
        converted = tmp_path + '_converted.png'
        if os.path.exists(converted):
            os.remove(converted)


@router.post("/parse-tasks", response_model=ParsedTasksResponse)
async def parse_tasks_from_text(request: ParseTasksRequest):
    """Use AI to parse tasks from free-form text with McKinsey-grade extraction"""

    api_key = os.getenv("ANTHROPIC_API_KEY")
    if not api_key:
        raise HTTPException(status_code=500, detail="ANTHROPIC_API_KEY not configured")

    client = Anthropic(api_key=api_key)

    prompt = f"""You are a senior McKinsey consultant specializing in workflow analysis and AI automation strategy.

USER INPUT:
{request.text}

Extract a structured, exhaustive task inventory from this input. Be a rigorous analyst:
- Decompose vague activities into atomic, measurable tasks
- Infer frequencies and time estimates from context clues
- Distinguish operational tasks from strategic ones
- Do not collapse distinct activities into one task

Respond ONLY with this exact JSON (no markdown, no code fences, no commentary):
{{
  "workflow_name": "Concise professional name (3-6 words)",
  "workflow_description": "One sharp sentence: who does what and why. Business-analyst tone.",
  "tasks": [
    {{
      "name": "Action-verb task name (e.g. 'Reconcile monthly expense reports')",
      "description": "What exactly happens, what inputs/outputs are involved, who is accountable",
      "frequency": "daily|weekly|monthly",
      "time_per_task": 30,
      "category": "data_entry|communication|analysis|creative|administrative|general",
      "complexity": "low|medium|high"
    }}
  ]
}}"""

    try:
        message = client.messages.create(
            model="claude-sonnet-4-20250514",
            max_tokens=3000,
            messages=[{"role": "user", "content": prompt}]
        )

        response_text = message.content[0].text.strip()
        if response_text.startswith('```'):
            lines = response_text.split('\n')
            response_text = '\n'.join(lines[1:-1])

        parsed = json.loads(response_text)
        return ParsedTasksResponse(**parsed)

    except Exception as e:
        print(f"Error parsing tasks: {e}")
        raise HTTPException(status_code=500, detail=f"Failed to parse tasks: {str(e)}")
